from flask import Flask, render_template, request, redirect, session
import sqlite3, math, os
from datetime import datetime,date
from werkzeug.security import generate_password_hash
from werkzeug.security import check_password_hash
from werkzeug.utils import secure_filename
from flask import flash, send_file
from google import genai
import json
import re
from markdown import markdown
from dotenv import load_dotenv
from pypdf import PdfReader
from helpers import (
    get_current_semester,
    get_task_stats,
    get_assignment_stats,
    get_attendance_stats,
    get_student_context,
    update_productivity_history,
    get_productivity_history,
    get_today_timetable,
    get_today_tasks,
    get_notes_context,
    get_today_assignments,
    get_today_information,
    remove_meal_times,
    calculate_free_slots,
    get_planner_context,
    get_note_text,
    get_summary,
    save_summary,
    get_ai_report,
    save_ai_report,
    get_study_plan,
    save_study_plan,
    extract_keywords,
    extract_chunk,
    retrieve_relevant_notes,
    get_cgpa
)

from ai_utils import (
    generate_ai_report,
    generate_study_plan,
    generate_note_summary,
    generate_weekly_review,
    generate_doubt_solver
)

app = Flask(__name__)
load_dotenv()

app.secret_key = os.getenv("SECRET_KEY")

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

conn = sqlite3.connect("database.db")
cursor = conn.cursor()

cursor.execute("""
CREATE TABLE IF NOT EXISTS users(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    username TEXT NOT NULL UNIQUE,
    password TEXT NOT NULL
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS tasks(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    task TEXT NOT NULL,
    status INTEGER DEFAULT 0,
    user_id INTEGER
)
""")

cursor.execute("""
    CREATE TABLE IF NOT EXISTS assignments(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        assignment_name TEXT NOT NULL,
        subject TEXT NOT NULL,
        due_date DATE NOT NULL,
        priority INTEGER NOT NULL,
        status INTEGER DEFAULT 0,
        user_id INTEGER,
        FOREIGN KEY(user_id) REFERENCES users(id)    
        )
""" )
cursor.execute("""
CREATE TABLE IF NOT EXISTS semesters(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    semester_no INTEGER NOT NULL,
    user_id INTEGER,
    FOREIGN KEY(user_id) REFERENCES users(id)
)
""")
cursor.execute("""
CREATE TABLE IF NOT EXISTS subjects(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    subject_name TEXT NOT NULL,
    credits INTEGER NOT NULL,
    grade INTEGER NOT NULL,
    semester_id INTEGER NOT NULL,
    user_id INTEGER,
    FOREIGN KEY(user_id) REFERENCES users(id),
    FOREIGN KEY(semester_id) REFERENCES semesters(id)
)
""")
cursor.execute("""
CREATE TABLE IF NOT EXISTS timetable(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    day TEXT NOT NULL,
    subject_name TEXT NOT NULL,
    start_time TEXT NOT NULL,
    end_time TEXT NOT NULL,
    semester_id INTEGER NOT NULL,
    user_id INTEGER,
    FOREIGN KEY(user_id) REFERENCES users(id),
    FOREIGN KEY(semester_id) REFERENCES semesters(id)
)
""")
cursor.execute("""
CREATE TABLE IF NOT EXISTS attendance(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    timetable_id INTEGER NOT NULL,
    date DATE NOT NULL,
    status INTEGER NOT NULL,
    user_id INTEGER,
    FOREIGN KEY(user_id) REFERENCES users(id),
    FOREIGN KEY(timetable_id) REFERENCES timetable(id)
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS productivity_history(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    score REAL,
    record_date DATE
)
""")


cursor.execute("""
CREATE TABLE IF NOT EXISTS notes(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    file_name TEXT NOT NULL,
    file_path TEXT NOT NULL,
    file_type TEXT NOT NULL,
    file_size INTEGER,
    subject TEXT,
    extracted_text TEXT,
    summary_json TEXT,
    upload_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    user_id INTEGER,
    FOREIGN KEY(user_id) REFERENCES users(id)
)
""")
cursor.execute("""
CREATE TABLE IF NOT EXISTS ai_reports(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER NOT NULL,
    report_type TEXT NOT NULL,
    report_json TEXT NOT NULL,
    generated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY(user_id) REFERENCES users(id)
)
""")
cursor.execute("""
    CREATE TABLE IF NOT EXISTS study_plans(
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER NOT NULL,
    plan_json TEXT NOT NULL,
    generated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY(user_id) REFERENCES users(id)
)
""")
conn.commit()
conn.close()

@app.route("/")
def landing():
    return render_template("landing.html")

@app.route("/tasks")
def tasks():
    if "username" in session:
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        filter_type=request.args.get("filter")
        search_text=request.args.get("search")
        query="SELECT * FROM tasks WHERE user_id=?"
        params=[session["user_id"]]
        if filter_type == "pending":
            query = query +" AND status=0"
        if filter_type == "completed":
            query = query + " AND status=1"
        if search_text:
            query = query + " AND task LIKE ?"
            params.append("%"+search_text+"%")
        cursor.execute(query,params)
        tasks=cursor.fetchall()
        message = None
        if search_text and len(tasks)==0:
            message = "No matching tasks found"
        elif len(tasks)==0:
            message = "No tasks yet. Add your first task!"
        cursor.execute(
            "SELECT COUNT(*) FROM tasks WHERE user_id=?",
            (session["user_id"],)
        )
        total_tasks = cursor.fetchone()[0]
        cursor.execute(
            "SELECT COUNT(*) FROM tasks WHERE user_id=? AND status=0",
            (session["user_id"],)
        )
        pending = cursor.fetchone()[0]
        cursor.execute(
            "SELECT COUNT(*) FROM tasks WHERE user_id=? AND status=1",
            (session["user_id"],)
        )
        completed = cursor.fetchone()[0]
        conn.close()
        return render_template("tasks.html",tasks=tasks,total_tasks=total_tasks,
                               pending=pending,completed=completed,message=message)

    else:
        return redirect("/login")

@app.route("/about")
def about():
    return "This is my first AI project"

@app.route("/add", methods=["POST"])
def add_task():
    if "username" in session:  
        task = request.form["task"]             # we store the text in variable and pass it
        conn = sqlite3.connect("database.db")   # first we are connecting to database
        cursor = conn.cursor()                  # using cursor we are able to apply CRUD 
                    
        cursor.execute(
            "INSERT INTO tasks(task, user_id) VALUES(?, ?)",
            (task, session["user_id"])
        )
        conn.commit() # now we are saving the changes
        conn.close()  # aftter saving the changes now we will close
        return redirect("/tasks")
    else:
        return redirect("/login")

@app.route("/delete/<int:id>",methods=["POST"])
def delete(id):
    if "username" in session:
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()  

        cursor.execute(
            "DELETE FROM tasks WHERE id=? AND user_id=?",
            (id, session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/tasks")  
    else:
        return redirect("/login")  # we are doing this so that even if someone who is not logged 
                                   #in tries to make changes to our project by typing webpage manually

@app.route("/update/<int:id>",methods=["GET","POST"])
def update(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    if request.method == "GET":
        cursor.execute(
            "SELECT * FROM tasks WHERE id=? AND user_id=?", 
            (id,session["user_id"])
        )
        task = cursor.fetchone()
        if task is None:   # This is so that one cant update others tasks and 
            conn.close()   #  it should reeturn to home page instead of crashing  
            return redirect("/tasks")
        conn.close()
        return render_template("update.html", task=task)
    else:
        new_task = request.form["task"]
        cursor.execute(
            "UPDATE tasks SET task=? WHERE id=?  AND user_id=?",
            (new_task, id,session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/tasks")

@app.route("/status/<int:id>",methods=["POST"])
def status(id):
    if "username" in session:
        conn = sqlite3.connect("database.db")   # THIS IS TO CHECK IF PROJECT IS COMPLETED OR PENDING
        cursor = conn.cursor()

        cursor.execute(
            " UPDATE tasks SET status=1 WHERE id=? AND user_id=?",
            (id,session["user_id"])
        )
        conn.commit()      
        conn.close()
        return redirect("/tasks")
    else:
        return redirect("/login")

@app.route("/register",methods=["GET","POST"])
def register():
    if "username" in session:
        return redirect("/dashboard")
    if request.method == "GET":
        return render_template("register.html",error=None)
    if request.method == "POST":
        user=request.form["username"]
        password=request.form["password"]
        confirm_password = request.form["confirm_password"]
        if password != confirm_password:
            return render_template( "register.html", error="Passwords do not match.")
        hashed_password = generate_password_hash(password)
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()

        cursor.execute(
        " SELECT * FROM users WHERE username=?",
        (user,)
        )
        existing_user = cursor.fetchone()
        if existing_user is None:
            cursor.execute(
                "INSERT INTO users(username,password) VALUES(?,?)",
            (user,hashed_password))
        else:
            conn.close()
            error="Username is Already Under Use"
            return render_template("register.html",error=error)
    conn.commit()      
    conn.close()
    return redirect("/login")
    
@app.route("/login",methods=["GET","POST"])
def login():
    if  "username" in session:
        return redirect("/dashboard")
    if request.method == "GET":
        return render_template("login.html",error=None)
    else:
        user=request.form["username"]
        password=request.form["password"]
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()

        cursor.execute(
            "SELECT * FROM users WHERE username=?",
            (user,)
        )
        existing_user = cursor.fetchone()
        if existing_user is None:
            conn.close()
            return render_template( "login.html", error="Invalid Credentials" )
        stored_hash = existing_user[2]
        if not check_password_hash(stored_hash, password):
            conn.close()
            return render_template( "login.html", error="Invalid Credentials")
        session["user_id"] = existing_user[0]
        session["username"] = existing_user[1]
        conn.close()
        return redirect("/dashboard")
    
@app.route("/logout")
def logout():
    session.clear()
    return redirect("/login")
    
@app.route("/assignments")
def assignments():
    if "username" in session:
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        cursor.execute(
            "SELECT * FROM assignments WHERE user_id=?",
            (session["user_id"],)
        )
        assig = cursor.fetchall()
        conn.close()
        return render_template("assignments.html",assig=assig)
    else:
        return redirect("/login")
    
@app.route("/add_assignments",methods=["POST"])
def add_assignments():
    if "username" in session:
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        add_assignment_name = request.form["assignment_name"]
        add_subject = request.form["subject"]
        add_due_date = request.form["due_date"]
        add_priority = int(request.form["priority"])
        cursor.execute(
            "INSERT INTO assignments(assignment_name, subject, due_date, priority, user_id)VALUES(?,?,?,?,?)", 
            (add_assignment_name, add_subject, add_due_date,add_priority, session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/assignments")
    else:
        return redirect("/login")

@app.route("/update_assignments/<int:id>",methods=["GET","POST"])
def update_assignments(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    if request.method == "GET":
        cursor.execute(
            "SELECT * FROM assignments WHERE id=? AND user_id=?", 
            (id,session["user_id"])
        )
        assig = cursor.fetchone()
        if assig is None:   # This is so that one cant update others tasks and 
            conn.close()   #  it should return to home page instead of crashing  
            return redirect("/")
        conn.close()
        return render_template("update_assignments.html", assig=assig)
    else:
        new_assignment_name = request.form["assignment_name"]
        new_subject = request.form["subject"]
        new_due_date = request.form["due_date"]
        new_priority = int(request.form["priority"])
        cursor.execute(
            "UPDATE assignments SET assignment_name=?, subject=?, due_date=?, priority=? WHERE id=? AND user_id=?",
            (new_assignment_name, new_subject, new_due_date, new_priority, id, session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/assignments") 

@app.route("/delete_assignments/<int:id>",methods=["POST"])
def delete_assignments(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()  
    cursor.execute(
        "DELETE FROM assignments WHERE id=? AND user_id=?",
        (id, session["user_id"])
    )
    conn.commit()
    conn.close()
    return redirect("/assignments") 

@app.route("/status_assignments/<int:id>",methods=["POST"])
def status_assignments(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db") 
    cursor = conn.cursor()
    cursor.execute(
        "UPDATE assignments SET status=1 WHERE id=? AND user_id=?",
        (id,session["user_id"])
    )
    conn.commit()      
    conn.close()
    return redirect("/assignments")
        
@app.route("/grades")
def grades():
    if "username" not in session:
        return redirect("/login") 
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT * FROM semesters WHERE user_id=?",
        (session["user_id"],)
    )
    grade = cursor.fetchall()
    semester_subjects={}
    semester_sgpa={}
    graph_data=[]
    for semester in grade:
        semester_id = semester[0]
        cursor.execute(
            "SELECT * FROM subjects WHERE semester_id=?",
            (semester_id,)
        )
        subjects=cursor.fetchall()
        semester_subjects[semester_id] = subjects # Now dictionary with sem_id and subject details
        total_points = 0
        total_credits = 0
        for subject in subjects:
            total_points = total_points + (subject[3]*subject[2]) # For SGPA 
            total_credits = total_credits + (subject[2])
        if total_credits > 0:
            sgpa = total_points / total_credits
        else:
            sgpa = 0 # THis is for a situation where user didnt add any subjects so "0/0" shouldn't happen
        semester_sgpa[semester_id] = round(sgpa, 2)
        graph_data.append(
            {
                "semester": f"Sem {semester[1]}",
                "sgpa": round(sgpa, 2)
            }
        ) # as it is dict we can use values right
    cgpa = get_cgpa(session["user_id"])
    conn.close()
    return render_template("grades.html",grade=grade,semester_subjects=semester_subjects, semester_sgpa=semester_sgpa, cgpa=cgpa,graph_data=graph_data)

@app.route("/add_semester",methods=["POST"])
def add_semester():
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    add_semester_no = request.form["semester_no"]
    cursor.execute(
        " SELECT * FROM semesters WHERE semester_no=? AND user_id=?",
        (add_semester_no,session["user_id"])
    )
    existing_semester_no=cursor.fetchone()
    cursor.execute(   # to make user select only semester 3 after semester 2
        "SELECT MAX(semester_no) FROM semesters WHERE user_id=?", 
        (session["user_id"],)
    )
    result = cursor.fetchone()
    if result[0] is None:
        max_semester = 0
    else:
        max_semester = result[0]

    if existing_semester_no: # checking for duplicates 
        cursor.execute(
            "SELECT * FROM semesters WHERE user_id=?",
            (session["user_id"],)
        )
        grade = cursor.fetchall()
        flash("Semester already exists.", "error")
        conn.close()
        return redirect("/grades")
    
    if int(add_semester_no) != max_semester + 1: # to allow only sem3 after sem2 using this
        cursor.execute(
            "SELECT * FROM semesters WHERE user_id=?",
            (session["user_id"],)
        )
        grade = cursor.fetchall()
        flash(f"Please add Semester {max_semester + 1} first.", "error")
        conn.close()
        return redirect("/grades")
    
    else:    
        cursor.execute(
            "INSERT INTO semesters(semester_no,user_id)VALUES(?,?)", 
            (add_semester_no, session["user_id"])
        )
    conn.commit()      
    conn.close()    
    return redirect("/grades")

@app.route("/add_subject/<int:semester_id>",methods=["GET","POST"])
def add_subject(semester_id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT id FROM semesters WHERE id=? AND user_id=?",
        (semester_id, session["user_id"])
    )
    semester = cursor.fetchone()
    if semester is None:
        conn.close()
        flash("Invalid semester.", "error")
        return redirect("/grades")
    if request.method == "GET":
        return render_template("add_subject.html",semester_id=semester_id)
    if request.method == "POST":
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        add_subject_name = request.form["subject_name"]
        add_grade = int(request.form["grade"])
        add_credits = int(request.form["credits"])
        cursor.execute(
            "INSERT INTO subjects(subject_name, grade, credits, user_id, semester_id)VALUES(?,?,?,?,?)", 
            (add_subject_name, add_grade, add_credits, session["user_id"],semester_id)
        )
        conn.commit()
        conn.close()
        return redirect("/grades")
    
@app.route("/delete_subject/<int:id>",methods=["POST"])
def delete_subjects(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()  
    cursor.execute(
        "DELETE FROM subjects WHERE id=? AND user_id=?",
        (id, session["user_id"])
    )
    conn.commit()
    conn.close()
    return redirect("/grades") 

@app.route("/update_subject/<int:id>",methods=["POST","GET"])
def update_subjects(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    if request.method == "GET":
        cursor.execute(
            "SELECT * FROM subjects WHERE id=? AND user_id=?", 
            (id,session["user_id"])
        )
        subject = cursor.fetchone()
        if subject is None:   
            conn.close()   
            return redirect("/grades")
        conn.close()
        return render_template("update_subject.html", subject=subject)
    else:
        new_subject_name = request.form["subject_name"]
        new_credits = int(request.form["credits"])
        new_grade = int(request.form["grade"])
        cursor.execute(
            "UPDATE subjects SET subject_name=?, credits=?, grade=? WHERE id=? AND user_id=?",
            (new_subject_name, new_credits, new_grade, id, session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/grades")
    
@app.route("/delete_semester/<int:id>",methods=["POST"])  
def delete_semester(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT semester_no FROM semesters WHERE id=? AND user_id=?", 
        (id,session["user_id"])
    )
    semester = cursor.fetchone()
    if semester is None:
        conn.close()
        flash("Semester not found.", "error")
        return redirect("/grades")
    deleting_semester = semester[0]
    cursor.execute(
        "SELECT MAX(semester_no) FROM semesters WHERE user_id=?", 
        (session["user_id"],)
    )
    maximum_semester = cursor.fetchone()[0]
    if maximum_semester == deleting_semester:
        cursor.execute(
            "DELETE FROM subjects WHERE user_id=? AND semester_id=?",
            (session["user_id"], id)
        )
        cursor.execute(
            "DELETE FROM semesters WHERE user_id=? AND id=?",
            (session["user_id"], id)
        )
        conn.commit()
        conn.close()
        flash("Semester deleted successfully.", "success")
        return redirect("/grades")
    else:
        conn.close()
        flash(f"Please Delete Semester { maximum_semester } first.", "error")
        return redirect("/grades")

@app.route("/confirm_delete_semester/<int:id>")
def confirm_delete(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT semester_no FROM semesters WHERE id=? AND user_id=?",
        (id, session["user_id"])
    )
    semester = cursor.fetchone()
    if semester is None:
        conn.close()
        return redirect("/grades")
    conn.close()
    return render_template("confirm_delete_semester.html",id=id,semester_no=semester[0])

@app.route("/timetable")
def timetable():
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        """SELECT * FROM semesters WHERE user_id=? ORDER BY semester_no DESC LIMIT 1 """,
        (session["user_id"],)
    )
    semester = cursor.fetchone()
    if not semester:
        message = "Please add a semester first from the Grades page."
        return render_template("timetable.html",semesters=None,message=message) 
    grouped_timetable = { 1: [], 2: [], 3: [], 4: [], 5: [] } # 1-> monday
    cursor.execute(
        "SELECT * FROM timetable WHERE semester_id=?",
        (semester[0],)
    )
    timetable_entries = cursor.fetchall()
    for entry in timetable_entries:
        day = int(entry[1])
        grouped_timetable[day].append(entry)         
    conn.close()
    return render_template("timetable.html", semester=semester,  semester_timetable=grouped_timetable)

@app.route("/add_class/<int:semester_id>",methods=["GET","POST"])
def add_class(semester_id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT id FROM semesters WHERE id=? AND user_id=?",
        (semester_id, session["user_id"])
    )
    semester = cursor.fetchone()
    if semester is None:
        conn.close()
        flash("Invalid semester.", "error")
        return redirect("/grades")
    if request.method == "GET":
        return render_template("add_class.html",semester_id=semester_id)
    if request.method == "POST":
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        subject = request.form["subject"]
        end_time = request.form["end_time"]
        start_time = request.form["start_time"]
        day=int(request.form["day"])
        if start_time >= end_time:
            error = "End time must be later than start time"
            return render_template("add_class.html",semester_id=semester_id,error=error)
        cursor.execute(
            """SELECT * FROM timetable WHERE semester_id=? AND day=? AND ( start_time < ? AND end_time > ?)""",
                (semester_id, day, end_time, start_time)
            )
        existing_slot = cursor.fetchone()
        if existing_slot:
            conn.close()
            error = "Time slot already occupied"
            return render_template("add_class.html",semester_id=semester_id,error=error)
        cursor.execute(
            "INSERT INTO timetable(subject_name, start_time, end_time, day, user_id, semester_id)VALUES(?,?,?,?,?,?)", 
            (subject, start_time, end_time, day, session["user_id"], semester_id)
        )
        conn.commit()
        conn.close()
        return redirect("/timetable")
    
@app.route("/delete_class/<int:id>", methods=["POST"])
def delete_class(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "DELETE FROM timetable WHERE id=? AND user_id=?",
        (id, session["user_id"])
    )
    conn.commit()
    conn.close()
    return redirect("/timetable")

@app.route("/update_class/<int:id>", methods=["GET", "POST"])
def update_class(id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    if request.method == "GET":
        cursor.execute(
            "SELECT * FROM timetable WHERE id=? AND user_id=?",
            (id, session["user_id"])
        )
        entry = cursor.fetchone()      
        if entry is None:
            conn.close()
            return redirect("/timetable")
        semester_id = entry[5]
        conn.close()
        return render_template("update_class.html",entry=entry)
    
    if request.method == "POST":
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        subject = request.form["subject"]
        end_time = request.form["end_time"]
        start_time = request.form["start_time"]
        day=int(request.form["day"])
        cursor.execute(
            "SELECT * FROM timetable WHERE id=? AND user_id=?",
            (id, session["user_id"])
        )
        entry = cursor.fetchone()
        semester_id = entry[5]
        if start_time >= end_time:
            error = "End time must be later than start time"
            return render_template("update_class.html",semester_id=semester_id,error=error,entry=entry)
        cursor.execute(
            """SELECT * FROM timetable WHERE semester_id=? AND day=? AND id!=? AND ( start_time < ? AND end_time > ?)""",
                (semester_id, day, id, end_time, start_time)
            )
        existing_slot = cursor.fetchone()
        if existing_slot:
            conn.close()
            error = "Time slot already occupied"
            return render_template("update_class.html",semester_id=semester_id,error=error,entry=entry)
        cursor.execute(
            "UPDATE timetable SET subject_name=?, day=?, start_time=?, end_time=? WHERE id=? AND user_id=?",
            (subject, day, start_time, end_time, id, session["user_id"])
        )
        conn.commit()
        conn.close()
        return redirect("/timetable")

@app.route("/dashboard")
def dashboard():
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT COUNT(*) FROM semesters WHERE user_id=?",
        (session["user_id"],)
    )
    semester_count = cursor.fetchone()[0]
    cursor.execute(
        "SELECT COUNT(*) FROM subjects WHERE user_id=?",
        (session["user_id"],)
    )
    subject_count = cursor.fetchone()[0]

    assignment_stats = get_assignment_stats(session["user_id"])
    pending_assignments = assignment_stats["pending_assignments"]

    task_stats = get_task_stats(session["user_id"])
    pending_tasks = task_stats["pending_tasks"]

    semester = get_current_semester(session["user_id"])
    if semester is None:
        return render_template("dashboard.html", no_semester=True)
    current_semester = semester["semester_id"]
    current_semester_no = semester["semester_no"]
    today_day = datetime.now().weekday() + 1
    if today_day>=6:
        today_classes_count=0
    else:
        cursor.execute(
            "SELECT COUNT(*) FROM timetable WHERE semester_id=? AND day=? AND user_id=?",
            (current_semester,today_day,session["user_id"])
        )
        today_classes_count = cursor.fetchone()[0]


    cgpa = get_cgpa(session["user_id"])
    
    attendance_stats = get_attendance_stats(session["user_id"],current_semester)

    attendance_alerts = [ 
                            { "subject": item["subject"], "percentage": item["percentage"] }
                            for item in attendance_stats["attendance_summary"] if item["percentage"] < 75       
                        ]

    notifications = []
    if pending_assignments == 1:
        notifications.append(
            "📌 You have 1 pending assignment"
        )
    elif pending_assignments > 1:
        notifications.append(
            f"📌 You have {pending_assignments} pending assignments"
        )
    if pending_tasks == 1:
        notifications.append( "📝 You have 1 pending task" )
    elif pending_tasks > 1:
        notifications.append( f"📝 You have {pending_tasks} pending tasks" )
    if today_classes_count == 1:
        notifications.append(   "📚 You have 1 class today" )
    elif today_classes_count > 1:
        notifications.append( f"📚 You have {today_classes_count} classes today" )
    for item in attendance_alerts:
        notifications.append( f"⚠ Attendance in {item['subject']} is {item['percentage']}%" )
    conn.close()
    return render_template("dashboard.html",pending_assignments=pending_assignments,pending_tasks=pending_tasks,today_classes_count=today_classes_count,cgpa=cgpa,current_semester_no=current_semester_no,attendance_alerts=attendance_alerts,
    notifications=notifications
    )

@app.route("/attendance",methods=["GET","POST"])
def attendance():
    if "username" not in session:
        return redirect("/login")
    if request.method == "GET":
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        cursor.execute(
            "SELECT * FROM semesters WHERE user_id=? ORDER BY semester_no DESC LIMIT 1", 
            (session["user_id"],)
        )
        existing_semester = cursor.fetchone()
        if existing_semester is None:
            return render_template("attendance.html",message="Please create a semester and timetable first.")
        current_semester = existing_semester[0]
        today_day = datetime.now().weekday() + 1
        message=None
        if today_day >=6:
            message = "No classes scheduled today."
        cursor.execute(
            "SELECT * FROM timetable WHERE semester_id=? AND day=? AND user_id=?",
            (current_semester,today_day,session["user_id"])
        )
        today_classes = cursor.fetchall()
        attendance_status = {}
        today_date = date.today()
        for class_item in today_classes:
            timetable_id = class_item[0]
            cursor.execute("""
                SELECT status FROM attendance WHERE timetable_id=? AND date=? AND user_id=?""",
                (timetable_id, today_date, session["user_id"])
            )
            record = cursor.fetchone()
            if record:
                attendance_status[timetable_id] = record[0]
        cursor.execute(
            """SELECT timetable.subject_name, SUM(CASE WHEN attendance.status = 1 THEN 1 ELSE 0 END) AS present,SUM(CASE WHEN attendance.status = 0 THEN 1  ELSE 0 END) AS absent, SUM(CASE WHEN attendance.status IN (0,1) THEN 1 ELSE 0 END) AS held
            FROM attendance
            JOIN timetable
            ON attendance.timetable_id = timetable.id
            WHERE attendance.user_id = ? AND timetable.semester_id = ?
            GROUP BY timetable.subject_name""",
             (session["user_id"], current_semester)
        )
        attendance_summary = cursor.fetchall()
        attendance_data = []
        for item in attendance_summary:
            subject_name = item[0]
            present = item[1]
            absent = item[2]
            held = item[3]
            if held > 0:
                percentage = round((present / held) * 100, 2)
            else:
                percentage = 0
            attendance_data.append({ "subject": subject_name, "present": present, "absent": absent, "held": held,"percentage": percentage})
        conn.close()
        return render_template("attendance.html",today_classes=today_classes, attendance_status=attendance_status,message=message, attendance_data=attendance_data)
    if request.method == "POST":
        conn = sqlite3.connect("database.db")
        cursor = conn.cursor()
        cursor.execute(
            "SELECT * FROM semesters WHERE user_id=? ORDER BY semester_no DESC LIMIT 1", 
            (session["user_id"],)
        )
        existing_semester = cursor.fetchone()
        if existing_semester is None:
            return render_template("attendance.html",message="Please create a semester and timetable first.")
        current_semester = existing_semester[0]
        message=None
        today_day = datetime.now().weekday() + 1
        today_date = date.today()
        if today_day >= 6:
            conn.close()
            return render_template("attendance.html",message="No classes scheduled today.",today_classes=[]
            )
        for timetable_id in request.form:
            status = request.form[timetable_id]
            cursor.execute(
                "SELECT * FROM attendance WHERE timetable_id=? AND date=? AND user_id=?",
                (timetable_id,today_date,session["user_id"])
            )
            entry = cursor.fetchone()
            if entry:
                cursor.execute(
                    "UPDATE attendance SET status=? WHERE timetable_id=? AND date=? AND user_id=?",
                    (status, timetable_id, today_date, session["user_id"])
                )
            else:
                cursor.execute(
                    "INSERT INTO attendance(timetable_id,date,status,user_id) VALUES(?,?,?,?)",
                   (timetable_id, today_date, status, session["user_id"])
                )
        conn.commit()
        conn.close()    
        return redirect("/attendance")

@app.route("/analytics")
def analytics():
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    penalty = 0
    reasons = []

    semester = get_current_semester(session["user_id"])
    if semester is None:
        return render_template( "attendance.html", message="Please create a semester and timetable first." )

    student_context = get_student_context(session["user_id"])

    update_productivity_history( session["user_id"], student_context["productivity_score"] )

    history = get_productivity_history( session["user_id"])
    
    task_stats = student_context["task_stats"]
    if task_stats["all_completed"]:
        reasons.append("All tasks completed")

    assignment_stats = student_context["assignment_stats"]
    if assignment_stats["overdue_count"] > 0:
        reasons.append( f'{assignment_stats["overdue_count"]} overdue assignments')
    else:
        reasons.append("No overdue assignments")
    
    attendance_stats = student_context["attendance_stats"]
    if attendance_stats["excellent_attendance"] :
        reasons.append("Excellent Attendance")
    else:
        for subject in attendance_stats["low_attendance_subjects"]:
            reasons.append( f"Attendance shortage in {subject}" )

    if task_stats["pending_tasks"] > 5:
        reasons.append( f'{task_stats["pending_tasks"]} pending tasks' )

    elif task_stats["pending_tasks"] > 0:
        reasons.append( f'{task_stats["pending_tasks"]} tasks remaining' )

    dates = history["dates"]
    scores = history["scores"]

    return render_template( "analytics.html",analytics_data=attendance_stats["attendance_summary"],assignment_analytics=assignment_stats["assignment_analytics"],total_tasks=task_stats["total_tasks"],pending_tasks=task_stats["pending_tasks"],completed_tasks=task_stats["completed_tasks"],completion_rate=task_stats["completion_rate"],task_message=
    task_stats["task_message"],productivity_score=student_context["productivity_score"],health=student_context["health"],
    reasons=reasons,dates=dates,scores=scores)

@app.route("/notes")
def notes():
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    search_text = request.args.get("search")
    subject_filter = request.args.get("subject")
    query = "SELECT * FROM notes WHERE user_id=?"
    params = [session["user_id"]]
    if search_text:
        query += """ AND ( file_name LIKE ? OR extracted_text LIKE ? )"""
        params.append("%" + search_text + "%")
        params.append("%" + search_text + "%")
    if subject_filter:
        query += " AND subject=?"        # to  show all files regarding that particular subject so basucally we are filtering
        params.append(subject_filter)    # we are ensuring that unit1,unit2 --> filenames will be shown if daa sub is pressed
    cursor.execute(query, params)
    notes = cursor.fetchall()
    ranked_notes = []
    for note in notes:
        found_in = ""
        if search_text:
            if search_text.lower() in note[1].lower():
                found_in = "File Name"
            elif search_text.lower() in note[6].lower():
                found_in = "Document Content"
        ranked_notes.append( (note, found_in))

    cursor.execute(
        "SELECT DISTINCT subject_name FROM timetable WHERE user_id=?", # in dropdown we are showing subjects right 
        (session["user_id"],)    # we are taking from timetable so automatically that sems subjects and prev sems should be shown
    )
    subjects = cursor.fetchall()
    conn.close()
    return render_template("notes.html",notes=ranked_notes,subjects=subjects)

@app.route("/upload", methods=["POST"])
def upload():
    if "username" not in session:
        return redirect("/login")
    uploaded_file = request.files["file"]
    subject = request.form["subject"]
    if uploaded_file.filename == "":   # if user didnt select file but clicked upload button then we are redirecting 
        return redirect("/notes")
    filename = secure_filename(uploaded_file.filename) # it is used to clear filenames
    base_name, extension = os.path.splitext(filename)
    file_type = extension.lower() # it is used to get filetype like pdf,ppt,.....
    allowed_types = [".pdf", ".pptx", ".txt", ".docx"]
    if file_type not in allowed_types:
        return redirect("/notes")
    upload_folder = "uploads/notes"
    new_filename = filename
    count = 1
    while os.path.exists(upload_folder + "/" + new_filename):
        new_filename = ( base_name + "(" + str(count) + ")" + extension ) # to get this way --> OperatingSystems(1).pdf
        count += 1     # we are adding count+=1 so that if once more OperatingSystem comes we get OperatingSystems(2).pdf
    file_path = os.path.join(upload_folder, new_filename)
    uploaded_file.save(file_path)
    file_size = os.path.getsize(file_path)
    extracted_text = ""
    if file_type == ".pdf":
        try:
            reader = PdfReader(file_path)    # opens the uploaded PDF
            for page in reader.pages:        # goes through every page.
                page_text = page.extract_text()   # page.extract_text()
                if page_text:
                    extracted_text += page_text + "\n"   # keeps adding all pages together.
        except Exception as e:
            print("PDF Extraction Error:", e)
    
    if file_type == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                extracted_text = f.read()
        except Exception as e:
            print("TXT Extraction Error:", e)

    if file_type == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        except Exception as e:
            print("DOCX Extraction Error:", e)

    if file_type == ".pptx":
        try:
            from pptx import Presentation
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):           #  checks if text exists and ignores images,blocks,diagrams...
                        extracted_text += shape.text + "\n"  
        except Exception as e:
            print("PPTX Extraction Error:", e)

    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        """ INSERT into notes(file_name,file_path,file_type,file_size,subject,extracted_text,user_id) 
        VALUES(?,?,?,?,?,?,?)
        """,
        (new_filename,file_path,file_type,file_size,subject,extracted_text,session["user_id"])
    )
    conn.commit()
    conn.close()
    return redirect("/notes")

@app.route("/open_note/<int:note_id>")
def open_note(note_id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT file_path FROM notes WHERE id=? AND user_id=?",
        (note_id,session["user_id"])
    )
    note = cursor.fetchone()
    if note is None:                    # Check whether this note belongs to the current user.
        flash("File not found")         # If note is None, either the note_id is invalid
        return redirect("/notes")       # or the user is trying to access someone else's file.
    file_path = note[0]
    if not os.path.exists(file_path):   # Get file path from database.
        flash("File does not exist.", "error")        # Then check whether the actual file still exists
        return redirect("/notes")       # inside uploads/notes folder.
    return send_file(file_path, as_attachment=False)

@app.route("/delete_note/<int:note_id>", methods=["POST"])
def delete_note(note_id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT file_path FROM notes WHERE id=? AND user_id=?",
        (note_id,session["user_id"])
    )
    note=cursor.fetchone()
    if note is None:
        conn.close()
        flash("File does not exist.", "error")
        return redirect("/notes")
    file_path = note[0]
    if os.path.exists(file_path):   # removing from uploads/motes
        os.remove(file_path)
    cursor.execute(
        "DELETE FROM notes WHERE id=? AND user_id=?",  # removing from database
        (note_id,session["user_id"])
    )
    conn.commit()
    conn.close()
    return redirect("/notes")

@app.route("/download_note/<int:note_id>")
def download_note(note_id):
    if "username" not in session:
        return redirect("/login")
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()
    cursor.execute(
        "SELECT file_path FROM notes WHERE id=? AND user_id=?",
        (note_id,session["user_id"])
    )
    note = cursor.fetchone()
    if note is None:                            # Check whether this note belongs to the current user.
        flash("File does not exist.", "error")  # If note is None, either the note_id is invalid
        return redirect("/notes")               # or the user is trying to access someone else's file.
    file_path = note[0]
    if not os.path.exists(file_path):            # Get file path from database.
        flash("File does not exist.", "error")   # Then check whether the actual file still exists
        return redirect("/notes")                # inside uploads/notes folder.
    return send_file(file_path, as_attachment=True)

@app.route("/ai_advisor")
def ai_advisor():
    if "username" not in session:
        return redirect("/login")
    saved = get_ai_report(session["user_id"], "advisor")
    semester = get_current_semester(session["user_id"])
    if semester is None:
        flash("Please create a semester first.", "warning")
        return redirect("/dashboard")

    if saved:
        report = json.loads(saved)
    else:
        try:
            student_context = get_student_context(session["user_id"])
            report = generate_ai_report(student_context)
            save_ai_report(session["user_id"],"advisor",report)
        except Exception as e:
            if "RESOURCE_EXHAUSTED" in str(e):
                flash("Daily AI quota reached. Please try again later.", "warning")
            else:
                flash("AI is currently busy. Please try again.", "warning")
            return redirect("/ai_advisor")    
    return render_template("ai_advisor.html",report=report)

@app.route("/regenerate_ai_advisor")
def regenerate_ai_advisor():
    if "username" not in session:
        return redirect("/login")
    semester = get_current_semester(session["user_id"])
    if semester is None:
        flash("Please create a semester first.", "warning")
        return redirect("/dashboard")
    
    try:
        student_context = get_student_context(session["user_id"])
        report = generate_ai_report(student_context)
        save_ai_report(session["user_id"], "advisor", report)
        flash("AI Advisor regenerated successfully.", "success")

    except Exception as e:
        if "RESOURCE_EXHAUSTED" in str(e):
            flash("Daily AI quota reached. Please try again later.", "warning")
        else:
            flash("AI is currently busy. Please try again.", "warning")

    return redirect("/ai_advisor")

@app.route("/ai_planner")
def ai_planner():
    if "username" not in session:
        return redirect("/login")
    saved = get_study_plan(session["user_id"])
    if saved:
        plan = json.loads(saved)
    else:
        try:
            student_context = get_planner_context(session["user_id"])
            plan = generate_study_plan(student_context)
            save_study_plan(session["user_id"],plan)
        except Exception as e:
            if "RESOURCE_EXHAUSTED" in str(e):
                flash("Daily AI quota reached. Please try again later.", "warning")
            else:
                flash("AI is currently busy. Please try again.", "warning")
            return redirect("/ai_planner") 
        
    return render_template("ai_planner.html",plan=plan)

@app.route("/regenerate_study_plan")
def regenerate_study_plan():
    if "username" not in session:
        return redirect("/login")

    try:
        student_context = get_planner_context(session["user_id"])
        plan = generate_study_plan(student_context)
        save_study_plan(session["user_id"], plan)
        flash("Study plan regenerated successfully.", "success")

    except Exception as e:
        if "RESOURCE_EXHAUSTED" in str(e):
            flash("Daily AI quota reached. Please try again later.", "warning")
        else:
            flash("AI is currently busy. Please try again.", "warning")

    return redirect("/ai_planner")

@app.route("/notes_summary/<int:note_id>")
def notes_summary(note_id):
    if "username" not in session:
        return redirect("/login")
    note = get_note_text(note_id, session["user_id"])
    if note is None:
        flash("This note could not be summarized.", "error")
        return redirect("/notes")
    summary_json = get_summary(note_id, session["user_id"])

    if summary_json:           # if we already have summary then we just show it, no use of extra tockens for every refresh
        summary = json.loads(summary_json)
    else:
        try:                                # if we generated a note_summary only then we use tockens
            summary = generate_note_summary(note["text"])
            save_summary(note_id,session["user_id"],summary)
        except Exception as e:                                           # if tockens are finished we show this
            if "RESOURCE_EXHAUSTED" in str(e):
                flash("Daily AI quota reached. Please try again later.", "warning")
            else:
                flash("AI is currently busy. Please try again.", "warning")
            return redirect("/notes")
    return render_template("notes_summary.html",note=note,summary=summary)

@app.route("/regenerate_summary/<int:note_id>")
def regenerate_summary(note_id):
    if "username" not in session:
        return redirect("/login")
    note = get_note_text(note_id, session["user_id"])
    if note is None:
        flash("Note not found.")
        return redirect("/notes")
    try:
        summary = generate_note_summary(note["text"])
        save_summary(note_id,session["user_id"],summary)
        flash("Summary regenerated successfully.", "success")
    except Exception as e:
        if "RESOURCE_EXHAUSTED" in str(e):
            flash("Daily AI quota reached. Please try again later.", "warning")
        else:
            flash("AI is currently busy. Please try again.", "warning")
        return redirect("/notes") 
    return redirect(f"/notes_summary/{note_id}")

@app.route("/weekly_review")
def weekly_review():
    if "username" not in session:
        return redirect("/login")
    semester = get_current_semester(session["user_id"])
    if semester is None:
        flash("Please create a semester first.", "warning")
        return redirect("/dashboard")
    saved = get_ai_report(session["user_id"],"weekly_review")
    if saved:
        review = json.loads(saved)
    else:
        try:
            student_context = get_student_context(session["user_id"])
            review = generate_weekly_review(student_context)
            save_ai_report(session["user_id"],"weekly_review",review)
        except Exception as e:
            if "RESOURCE_EXHAUSTED" in str(e):
                flash("Daily AI quota reached. Please try again later.", "warning")
            else:
                flash("AI is currently busy. Please try again.", "warning")
            return redirect("/weekly_review") 
    return render_template("weekly_review.html",review=review)

@app.route("/regenerate_weekly_review")
def regenerate_weekly_review():
    if "username" not in session:
        return redirect("/login")
    semester = get_current_semester(session["user_id"])
    if semester is None:
        flash("Please create a semester first.", "warning")
        return redirect("/dashboard")

    try:
        student_context = get_student_context(session["user_id"])
        review = generate_weekly_review(student_context)
        save_ai_report(session["user_id"], "weekly_review", review)
        flash("Weekly review regenerated successfully.", "success")

    except Exception as e:
        if "RESOURCE_EXHAUSTED" in str(e):
            flash("Daily AI quota reached. Please try again later.", "warning")
        else:
            flash("AI is currently busy. Please try again.", "warning")

    return redirect("/weekly_review")

@app.route("/doubt_solver", methods=["GET", "POST"])
def doubt_solver():
    if "username" not in session:
        return redirect("/login")
    answer = None
    sources = []
    question = ""
    answer = None
    summary = None
    confidence = None
    if request.method == "POST":
        question = request.form.get("question", "").strip()
        if question:
            matched_notes = retrieve_relevant_notes( session["user_id"], question)
            if matched_notes:
                context = ""
                for note in matched_notes:
                    context += (
                        f"Source: {note['file_name']}\n\n"
                        f"{note['text']}\n\n"                       # we are creating filename below it extracted text 
                        "----------------------------------\n\n"  # and below it we are drawing a line to seperate another file and text
                    )
                try:
                    response = generate_doubt_solver(question, context)
                    answer = markdown(response["answer"].strip(), extensions=["extra"])
                    summary = response["summary"].strip()
                    confidence = response["confidence"].strip()
                    sources = [note["file_name"] for note in matched_notes]
                except Exception as e:
                    if "RESOURCE_EXHAUSTED" in str(e):
                        flash("Daily AI quota reached. Please try again later.", "warning")
                    else:
                        flash("AI is currently busy. Please try again.", "warning")
                    return redirect("/doubt_solver") 
            else:
                answer = (
                    "I couldn't find the answer in your uploaded notes."
                )
    
    return render_template( "doubt_solver.html", answer=answer, summary=summary, confidence=confidence, sources=sources, question=question )



if __name__ == "__main__":
    app.run(debug=False)